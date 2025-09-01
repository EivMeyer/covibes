#!/usr/bin/env python3
"""
Create ColabVibe PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    slides_data = [
        ("ColabVibe", "The AI Workforce Platform\n\nFrom single-player AI assistants to a multi-player AI workforce.\n\n(Tromsø, September 2025)"),
        ("The Problem", "AI Development is Trapped in a Single-User Paradigm\n\nToday's AI is a personal tool, not a team member. This creates bottlenecks and stops progress when the developer logs off."),
        ("The Solution", "The First Team-Wide AI Workforce Platform\n\nCore Features:\n• Team-Wide Spawning: Anyone can deploy AI agents\n• Shared Context: Agents access the full workspace\n• Mobile-First Orchestration: Manage from anywhere"),
        ("The Paradigm Shift", "Software is Moving From MS Word to Google Docs\n\nPast: Grammar checker in a private .doc file (helps one person).\nFuture: Shared, live workspace where AI agents are collaborators.\n\nWe are doing for AI development what Google Docs did for documents."),
        ("How It Works", "Fixing a Critical Bug: Hours vs. Minutes\n\nToday (AI-Assisted):\n👨‍💻 Human-driven (Dev at desk)\n💡 AI Assists (Autocomplete)\n⏳ Sequential & Slow\n\nWith ColabVibe (AI Workforce):\n📱 Human-Directed (PM on phone)\n🤖 AI Acts (Autonomous fix)\n⚡ Parallel & Fast"),
        ("The Product", "Multi-Agent Collaboration in Action\n\nHighlights from Demo:\n• Parallel Work: Agent A updates frontend while Agent B updates backend\n• Real-Time Visibility: Team sees live progress\n• Human-in-the-Loop: Developer approves deployment\n\nWatch the demo: www.youtube.com/your-video-link"),
        ("Market Opportunity", "Creating a New Category: AI Workforce Management\n\n$26B+ Developer Tools Market\n100M+ Developers Worldwide\n92% of developers use AI, but in isolation.\n\nTrend: Remote work demands collaborative, always-on development."),
        ("Go-to-Market Strategy", "Our Plan for the First 1,000 Users\n\nPhase 1 (Alpha): 15 Norwegian tech companies as design partners\nPhase 2 (Beta): Product Hunt, Hacker News, technical content\nPhase 3 (Public): Developer communities & referrals"),
        ("Business Model", "AI Workforce-as-a-Service (WaaS)\n\nPricing Tiers:\n• Starter (small teams)\n• Professional (growing businesses)\n• Enterprise (large-scale needs)\n\nModel: Simple per-user, per-month subscription"),
        ("Competitive Advantage", "Why We Win: Built for Teams\n\n• Multi-Agent Coordination: Competitors are single-player\n• Shared Workspace Context: Our agents have full access\n• Mobile-First Orchestration: Competitors are desktop-only"),
        ("The Team", "The Team to Build the Future of Development\n\n[Add founder headshots, names, titles, and 2-3 accomplishments each]"),
        ("Vision & Roadmap", "Vision: AI Teams That Work While You Sleep\n\nNow: Any team member can spawn AI agents\n12 Months: Agents collaborate with each other\n24 Months+: Fully autonomous AI teams managed by human orchestration"),
        ("The Ask & Use of Funds", "Join Us in Building the Future\n\nFunding Ask: We are raising $[Amount] to achieve 12-month milestones\n\nUse of Funds:\n• 60% Product & Engineering (hire 3 developers in Tromsø)\n• 25% Go-to-Market (hire 1 marketer in Norway)\n• 15% Operations"),
        ("Contact", "ColabVibe\n\nReady for a Live Demo?\n\nContact: [Your Name], [Your Email], [colabvibe.com]\n\nThe future of development isn't just AI-assisted—it's AI-workforce managed.")
    ]
    
    for title, content in slides_data:
        # Use blank layout for better control over positioning
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add and position title manually
        left = Inches(0.5)
        top = Inches(0.3) if title == "ColabVibe" else Inches(0.2)
        width = Inches(9)
        height = Inches(1.2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.clear()
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.alignment = PP_ALIGN.CENTER
        title_font = title_para.font
        title_font.size = Pt(44) if title == "ColabVibe" else Pt(36)
        title_font.bold = True
        title_font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        
        # Add content with reduced gap
        content_top = Inches(1.8) if title == "ColabVibe" else Inches(1.2)
        content_height = Inches(3.5)
        
        content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
        text_frame = content_box.text_frame
        text_frame.clear()
        
        # Add content with proper formatting
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(20) if title == "ColabVibe" else Pt(18)
        p.alignment = PP_ALIGN.LEFT if "•" in content else PP_ALIGN.CENTER
        p.space_after = Pt(6)
        
        # Center content vertically in its box
        text_frame.vertical_anchor = MSO_ANCHOR.TOP if "•" in content else MSO_ANCHOR.MIDDLE
        
        # Special formatting for title slide
        if title == "ColabVibe":
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Save presentation
    prs.save('ColabVibe_Presentation.pptx')
    print("✅ PowerPoint presentation created: ColabVibe_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()